from flask import Flask, request, jsonify, send_file, send_from_directory
from flask_cors import CORS
import json
import os
import tempfile
from datetime import datetime
import openai
import anthropic
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import logging

app = Flask(__name__)
CORS(app)

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class LLMOrchestrator:
    def __init__(self, provider, api_key):
        self.provider = provider.lower()
        self.api_key = api_key
        
    def generate_outline_with_notes(self, text, guidance="", include_notes=False):
        """Generate slide outline with optional speaker notes"""
        notes_instruction = """
        Also generate speaker notes for each slide. Include detailed explanations, examples, and talking points.
        """ if include_notes else ""
        
        points_format = '"points": ["Point 1", "Point 2"], "notes": "Detailed speaker notes for this slide"' if include_notes else '"points": ["Point 1", "Point 2"]'
        
        prompt = f"""
        Convert the following text into a PowerPoint presentation outline. 
        ,heading should not have slide numbers etc,.
        Create a JSON structure with slides containing titles and bullet points.
        {notes_instruction}
        
        Guidelines: {guidance if guidance else "Create a professional presentation with clear structure "}
        
        Text to convert:
        {text}
        
        Return ONLY a JSON object in this format:
        {{
            "title": "Presentation Title",
            "slides": [
                {{"title": "Slide Title", {points_format}}},
                {{"title": "Next Slide", {points_format}}}
            ]
        }}
        """
        
        try:
            if self.provider == "openai":
                client = openai.OpenAI(api_key=self.api_key)
                response = client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.7
                )
                content = response.choices[0].message.content
                
            elif self.provider == "anthropic":
                client = anthropic.Anthropic(api_key=self.api_key)
                response = client.messages.create(
                    model="claude-3-sonnet-20240229",
                    max_tokens=2000,
                    messages=[{"role": "user", "content": prompt}]
                )
                content = response.content[0].text
                
            elif self.provider == "groq":
                client = Groq(api_key=self.api_key)
                response = client.chat.completions.create(
                    model="llama-3.1-8b-instant",
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.7
                )
                content = response.choices[0].message.content
                
            else:
                raise ValueError(f"Unsupported provider: {self.provider}")
            
            # Clean and parse JSON
            content = content.strip()
            if content.startswith("```json"):
                content = content[7:]
            if content.endswith("```"):
                content = content[:-3]
            
            return json.loads(content)
            
        except Exception as e:
            logger.error(f"LLM generation error: {str(e)}")
            # Fallback outline
            fallback_slides = [
                {"title": "Introduction", "points": ["Overview of the topic", "Key objectives"]},
                {"title": "Main Content", "points": ["Key point 1", "Key point 2", "Key point 3"]},
                {"title": "Conclusion", "points": ["Summary", "Next steps"]}
            ]
            
            if include_notes:
                for slide in fallback_slides:
                    slide["notes"] = f"Speaker notes for {slide['title']}: Expand on the key points and provide examples."
            
            return {
                "title": "Generated Presentation",
                "slides": fallback_slides
            }

class PPTGenerator:
    def __init__(self, template_path=None):
        self.presentation = None
        self.template_path = template_path
        self.content_layouts = []  # List of suitable content layouts
        self.current_layout_index = 0  # Index for rotating through layouts
        
    def create_presentation(self, outline_data):
        """Create PowerPoint presentation from outline data"""
        try:
            # Use template if provided, otherwise create new presentation
            if self.template_path and os.path.exists(self.template_path):
                logger.info(f"Using template: {self.template_path}")
                self.presentation = Presentation(self.template_path)
                logger.info(f"Template loaded with {len(self.presentation.slide_layouts)} layouts")
                logger.info(f"Template has {len(self.presentation.slides)} existing slides")
                
                # Log available layouts for debugging
                for i, layout in enumerate(self.presentation.slide_layouts):
                    logger.info(f"Layout {i}: {layout.name}")
                
                # Identify and catalog content layouts for variety
                self._identify_content_layouts()
                
                # Clear existing slides but preserve the layouts and master slides
                # This keeps the template design but removes content
                while len(self.presentation.slides) > 0:
                    rId = self.presentation.slides._sldIdLst[0].rId
                    self.presentation.part.drop_rel(rId)
                    del self.presentation.slides._sldIdLst[0]
                
                logger.info("Cleared existing slides, keeping layouts and design")
            else:
                logger.info("Creating new presentation with default template")
                self.presentation = Presentation()
            
            # Create title slide using the first layout (typically title slide layout)
            title_layout = self.presentation.slide_layouts[0]
            title_slide = self.presentation.slides.add_slide(title_layout)
            logger.info(f"Created title slide with layout: {title_layout.name}")
            
            # Set title content
            title_shape = None
            subtitle_shape = None
            
            for shape in title_slide.shapes:
                if shape.is_placeholder:
                    if shape.placeholder_format.idx == 0:  # Title placeholder
                        title_shape = shape
                    elif shape.placeholder_format.idx == 1:  # Subtitle placeholder
                        subtitle_shape = shape
            
            # Also try the shapes.title method as fallback
            if not title_shape and hasattr(title_slide.shapes, 'title'):
                title_shape = title_slide.shapes.title
            
            if title_shape:
                if hasattr(title_shape, 'text_frame'):
                    title_shape.text_frame.text = outline_data.get("title", "Generated Presentation")
                elif hasattr(title_shape, 'text'):
                    title_shape.text = outline_data.get("title", "Generated Presentation")
                logger.info(f"Set title: {outline_data.get('title', 'Generated Presentation')}")
            
            if subtitle_shape:
                if hasattr(subtitle_shape, 'text_frame'):
                    subtitle_shape.text_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
                elif hasattr(subtitle_shape, 'text'):
                    subtitle_shape.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
                logger.info("Set subtitle with date")
            
            # Content slides
            for slide_data in outline_data.get("slides", []):
                self._create_content_slide(slide_data)
                
            return True
            
        except Exception as e:
            logger.error(f"PPT generation error: {str(e)}")
            return False
    
    def _identify_content_layouts(self):
        """Identify and catalog all suitable content layouts from the template"""
        self.content_layouts = []
        
        # Keywords that indicate content layouts
        content_keywords = [
            'content', 'bullet', 'text', 'two content', 'comparison', 
            'title and content', 'section header', 'picture with caption',
            'content with caption', 'blank'
        ]
        
        title_keywords = ['title', 'section']
        
        for i, layout in enumerate(self.presentation.slide_layouts):
            layout_name = layout.name.lower()
            
            # Skip pure title layouts (we handle title slide separately)
            if any(keyword in layout_name for keyword in ['title slide', 'title only']):
                continue
            
            # Include layouts that are good for content
            is_content_layout = False
            
            # Check by name
            if any(keyword in layout_name for keyword in content_keywords):
                is_content_layout = True
            
            # Check by layout structure (has content placeholders)
            elif self._has_content_placeholders(layout):
                is_content_layout = True
            
            # Include layouts that aren't purely title layouts
            elif not any(keyword in layout_name for keyword in title_keywords):
                is_content_layout = True
            
            if is_content_layout:
                self.content_layouts.append((i, layout))
                logger.info(f"Added content layout {i}: {layout.name}")
        
        # If no content layouts found, use default layouts
        if not self.content_layouts:
            for i in range(min(len(self.presentation.slide_layouts), 3)):
                if i > 0:  # Skip the first one (usually title slide)
                    self.content_layouts.append((i, self.presentation.slide_layouts[i]))
            logger.info("No content layouts found by name, using default layouts")
        
        logger.info(f"Identified {len(self.content_layouts)} content layouts for variety")
    
    def _has_content_placeholders(self, layout):
        """Check if a layout has content placeholders (not just title)"""
        try:
            # Create a temporary slide to check placeholders
            temp_slide = self.presentation.slides.add_slide(layout)
            has_content = False
            content_placeholder_count = 0
            
            for placeholder in temp_slide.placeholders:
                # Look for non-title placeholders that can hold text
                if placeholder.placeholder_format.idx > 0:
                    if hasattr(placeholder, 'text_frame'):
                        content_placeholder_count += 1
                        has_content = True
            
            # Also check for text boxes or shapes that can hold content
            if not has_content:
                for shape in temp_slide.shapes:
                    if hasattr(shape, 'text_frame') and not shape.is_placeholder:
                        has_content = True
                        break
            
            # Remove the temporary slide safely
            try:
                rId = self.presentation.slides._sldIdLst[-1].rId
                self.presentation.part.drop_rel(rId)
                del self.presentation.slides._sldIdLst[-1]
            except:
                pass  # If removal fails, continue anyway
            
            logger.info(f"Layout '{layout.name}' has {content_placeholder_count} content placeholders, suitable: {has_content}")
            return has_content
        except Exception as e:
            logger.warning(f"Error checking layout '{layout.name}': {str(e)}")
            return True  # Assume it's usable if we can't check
    
    def _get_next_content_layout(self):
        """Get the next content layout for variety"""
        if not self.content_layouts:
            # Fallback to standard layouts
            if len(self.presentation.slide_layouts) > 1:
                return self.presentation.slide_layouts[1]
            else:
                return self.presentation.slide_layouts[0]
        
        # Rotate through available content layouts
        layout_index, layout = self.content_layouts[self.current_layout_index]
        self.current_layout_index = (self.current_layout_index + 1) % len(self.content_layouts)
        
        logger.info(f"Using content layout {layout_index}: {layout.name} (rotation {self.current_layout_index}/{len(self.content_layouts)})")
        return layout
    
    def _get_smart_layout(self, slide_data):
        """Intelligently select layout based on slide content"""
        if not self.content_layouts:
            return self._get_next_content_layout()
        
        slide_title = slide_data.get("title", "").lower()
        points = slide_data.get("points", [])
        num_points = len(points)
        
        # Try to match content type with appropriate layout
        suitable_layouts = []
        
        for layout_index, layout in self.content_layouts:
            layout_name = layout.name.lower()
            
            # Match based on content characteristics
            if num_points <= 2 and 'two content' in layout_name:
                suitable_layouts.append((layout_index, layout, 3))  # High priority
            elif num_points > 4 and 'bullet' in layout_name:
                suitable_layouts.append((layout_index, layout, 3))  # High priority
            elif 'comparison' in layout_name and ('vs' in slide_title or 'comparison' in slide_title):
                suitable_layouts.append((layout_index, layout, 4))  # Very high priority
            elif 'section' in layout_name and ('introduction' in slide_title or 'conclusion' in slide_title):
                suitable_layouts.append((layout_index, layout, 2))  # Medium priority
            else:
                suitable_layouts.append((layout_index, layout, 1))  # Default priority
        
        if suitable_layouts:
            # Sort by priority and pick the best match
            suitable_layouts.sort(key=lambda x: x[2], reverse=True)
            best_layout = suitable_layouts[0][1]
            logger.info(f"Smart layout selection: {best_layout.name} for slide '{slide_title}'")
            return best_layout
        
        # Fallback to rotation
        return self._get_next_content_layout()
    
    def _create_content_slide(self, slide_data):
        """Create a content slide with title, bullet points, and optional speaker notes"""
        # Use smart layout selection for variety and content matching
        content_layout = self._get_smart_layout(slide_data)
        
        # Create the slide using the selected template layout
        slide = self.presentation.slides.add_slide(content_layout)
        logger.info(f"Created slide '{slide_data.get('title', 'Slide Title')}' with layout: {content_layout.name}")
        
        # Log all placeholders for debugging
        for i, placeholder in enumerate(slide.placeholders):
            logger.info(f"Placeholder {i}: idx={placeholder.placeholder_format.idx}, name={getattr(placeholder, 'name', 'no name')}")
        
        # Set title - try multiple methods to ensure it works
        title_set = False
        
        # Method 1: Use placeholder index 0 (standard title placeholder)
        try:
            title_placeholder = slide.placeholders[0]
            if hasattr(title_placeholder, 'text_frame'):
                title_placeholder.text_frame.text = slide_data.get("title", "Slide Title")
            elif hasattr(title_placeholder, 'text'):
                title_placeholder.text = slide_data.get("title", "Slide Title")
            title_set = True
            logger.info(f"Set title via placeholder[0]: {slide_data.get('title', 'Slide Title')}")
        except:
            pass
        
        # Method 2: Use shapes.title if available
        if not title_set:
            try:
                if hasattr(slide.shapes, 'title') and slide.shapes.title:
                    slide.shapes.title.text = slide_data.get("title", "Slide Title")
                    title_set = True
                    logger.info(f"Set title via shapes.title: {slide_data.get('title', 'Slide Title')}")
            except:
                pass
        
        # Method 3: Find title placeholder manually
        if not title_set:
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    try:
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.text = slide_data.get("title", "Slide Title")
                        elif hasattr(shape, 'text'):
                            shape.text = slide_data.get("title", "Slide Title")
                        title_set = True
                        logger.info(f"Set title via manual search: {slide_data.get('title', 'Slide Title')}")
                        break
                    except:
                        continue
        
        # Add content - try multiple methods in order of preference
        content_added = False
        points = slide_data.get("points", [])
        logger.info(f"Adding {len(points)} bullet points")
        
        if points:  # Only try to add content if there are points
            # Method 1: Look for content placeholders, prioritizing text content
            content_placeholder = None
            
            # Strategy 1: Look for the most common content placeholder (idx=1)
            try:
                for shape in slide.shapes:
                    if (shape.is_placeholder and 
                        shape.placeholder_format.idx == 1 and
                        hasattr(shape, 'text_frame')):
                        content_placeholder = shape
                        logger.info(f"Found standard content placeholder idx=1 (name: {getattr(shape, 'name', 'no name')})")
                        break
            except:
                pass
            
            # Strategy 2: Look for any content placeholder with "content" in the name
            if not content_placeholder:
                for shape in slide.shapes:
                    if (shape.is_placeholder and 
                        shape.placeholder_format.idx > 0 and
                        hasattr(shape, 'text_frame')):
                        shape_name = getattr(shape, 'name', '').lower()
                        if 'content' in shape_name and 'picture' not in shape_name:
                            content_placeholder = shape
                            logger.info(f"Found named content placeholder idx={shape.placeholder_format.idx} (name: {getattr(shape, 'name', 'no name')})")
                            break
            
            # Strategy 3: Look for any non-picture placeholder that can hold text
            if not content_placeholder:
                for shape in slide.shapes:
                    if (shape.is_placeholder and 
                        shape.placeholder_format.idx > 0 and
                        hasattr(shape, 'text_frame')):
                        shape_name = getattr(shape, 'name', '').lower()
                        # Skip obvious picture placeholders
                        if not any(word in shape_name for word in ['picture', 'image', 'photo', 'chart', 'table']):
                            content_placeholder = shape
                            logger.info(f"Found text placeholder idx={shape.placeholder_format.idx} (name: {getattr(shape, 'name', 'no name')})")
                            break
            
            # Strategy 4: Use any placeholder as last resort (even picture placeholders)
            if not content_placeholder:
                for shape in slide.shapes:
                    if (shape.is_placeholder and 
                        shape.placeholder_format.idx > 0 and
                        hasattr(shape, 'text_frame')):
                        content_placeholder = shape
                        logger.info(f"Using fallback placeholder idx={shape.placeholder_format.idx} (name: {getattr(shape, 'name', 'no name')})")
                        break
            
            # Method 2: Try to add content to the found placeholder
            if content_placeholder:
                try:
                    text_frame = content_placeholder.text_frame
                    logger.info(f"Text frame found, clearing and adding {len(points)} points")
                    text_frame.clear()
                    
                    for i, point in enumerate(points):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = point
                        p.level = 0
                        # Try to set font size if possible
                        try:
                            p.font.size = Pt(18)
                        except:
                            pass
                    
                    content_added = True
                    logger.info(f"✅ Successfully added content via placeholder (idx {content_placeholder.placeholder_format.idx})")
                    
                    # Verify content was actually added
                    final_text = text_frame.text.strip()
                    if final_text:
                        logger.info(f"✅ Content verification: Text frame now contains: {final_text[:50]}...")
                    else:
                        logger.warning(f"⚠️ Content verification failed: Text frame appears empty")
                        
                except Exception as e:
                    logger.error(f"❌ Failed to add content to placeholder: {str(e)}")
                    content_placeholder = None  # Reset so we can try other methods
            
            # Method 3: Find any text frame that's not a title (more aggressive)
            if not content_added:
                logger.info("Trying alternative text frame method...")
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        try:
                            # Check if this might be a title by comparing with slide title
                            current_text = shape.text_frame.text.strip()
                            slide_title = slide_data.get("title", "").strip()
                            
                            # Skip if it's likely the title shape
                            if current_text == slide_title:
                                logger.info(f"Skipping title shape: {current_text}")
                                continue
                            
                            # Try to use this text frame
                            text_frame = shape.text_frame
                            logger.info(f"Found alternative text frame, current content: '{current_text[:30]}...'")
                            
                            # Clear only if it's empty or has minimal content
                            if len(current_text) < 10:
                                text_frame.clear()
                                start_paragraph = text_frame.paragraphs[0]
                            else:
                                # Add content after existing content
                                start_paragraph = text_frame.add_paragraph()
                            
                            # Add our content
                            for i, point in enumerate(points):
                                if i == 0:
                                    p = start_paragraph
                                else:
                                    p = text_frame.add_paragraph()
                                
                                p.text = f"• {point}"
                                p.level = 0
                                try:
                                    p.font.size = Pt(18)
                                except:
                                    pass
                            
                            content_added = True
                            logger.info(f"✅ Added content via alternative text frame")
                            
                            # Verify content
                            final_text = text_frame.text.strip()
                            logger.info(f"✅ Final text frame content: {final_text[:100]}...")
                            break
                            
                        except Exception as e:
                            logger.warning(f"Could not use text frame: {str(e)}")
                            continue
            
            # Method 4: Add a new text box if no suitable placeholder found
            if not content_added:
                try:
                    from pptx.shapes.textbox import TextBox
                    from pptx.util import Inches
                    
                    # Add a text box in a reasonable position
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(4)
                    
                    textbox = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = textbox.text_frame
                    text_frame.clear()
                    
                    for i, point in enumerate(points):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = f"• {point}"
                        p.level = 0
                        try:
                            p.font.size = Pt(18)
                        except:
                            pass
                    
                    content_added = True
                    logger.info(f"Added content via new text box")
                except Exception as e:
                    logger.error(f"Could not create text box: {str(e)}")
        
        # Log the result
        if content_added:
            logger.info(f"Successfully added content to slide: {slide_data.get('title', 'Slide Title')}")
        elif points:
            logger.warning(f"Created slide with title only: {slide_data.get('title', 'Slide Title')}")
        else:
            logger.info(f"Created title-only slide (no content points): {slide_data.get('title', 'Slide Title')}")
        
        # Add speaker notes if provided
        notes_text = slide_data.get("notes", "")
        if notes_text:
            try:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = notes_text
                logger.info(f"Added speaker notes")
            except Exception as e:
                logger.error(f"Error adding speaker notes: {str(e)}")
        
        if title_set and content_added:
            logger.info(f"Successfully created slide: {slide_data.get('title', 'Slide Title')}")
        elif title_set:
            logger.warning(f"Created slide with title only: {slide_data.get('title', 'Slide Title')}")
        else:
            logger.error(f"Failed to create slide properly: {slide_data.get('title', 'Slide Title')}")
        
        return slide
    
    def save_presentation(self, filename):
        """Save presentation to file"""
        if self.presentation:
            self.presentation.save(filename)
            return True
        return False

@app.route('/', methods=['GET'])
def home():
    """Serve the main HTML page"""
    try:
        # Get the path to the frontend directory
        return send_from_directory('.', 'index.html')
    except Exception as e:
        logger.error(f"Error serving home page: {str(e)}")
        return jsonify({"error": "Frontend not found"}), 404

@app.route('/test', methods=['GET'])
def test_page():
    """Serve test page to verify server is working"""
    return send_from_directory('.', 'test.html')

@app.route('/health', methods=['GET'])
def health_check():
    return jsonify({"status": "healthy", "timestamp": datetime.now().isoformat()})

@app.route('/generate', methods=['POST'])
def generate_presentation():
    try:
        # Handle both form data (with file) and JSON data
        if request.content_type and 'multipart/form-data' in request.content_type:
            # Form data with potential file upload
            text = request.form.get('text', '').strip()
            provider = request.form.get('provider', '').strip().lower()
            api_key = request.form.get('api_key', '').strip()
            guidance = request.form.get('guidance', '').strip()
            include_notes = request.form.get('include_notes', 'false').lower() == 'true'
            
            # Handle template file upload
            template_file = request.files.get('template')
            template_path = None
            
            if template_file and template_file.filename:
                if not template_file.filename.lower().endswith(('.pptx', '.potx')):
                    return jsonify({"error": "Template must be a .pptx or .potx file"}), 400
                
                # Save template to temporary location
                temp_dir = tempfile.mkdtemp()
                template_path = os.path.join(temp_dir, "template.pptx")
                template_file.save(template_path)
                logger.info(f"Template uploaded and saved to: {template_path}")
        else:
            # JSON data
            data = request.get_json()
            if not data:
                return jsonify({"error": "No data provided"}), 400
            
            text = data.get('text', '').strip()
            provider = data.get('provider', '').strip().lower()
            api_key = data.get('api_key', '').strip()
            guidance = data.get('guidance', '').strip()
            include_notes = data.get('include_notes', False)
            template_path = None
        
        # Validation
        if not text:
            return jsonify({"error": "Text content is required"}), 400
        
        if not provider or provider not in ['openai', 'anthropic', 'groq']:
            return jsonify({"error": "Valid provider (openai, anthropic, groq) is required"}), 400
        
        if not api_key:
            return jsonify({"error": "API key is required"}), 400
        
        logger.info(f"Generating presentation with {provider} provider, notes: {include_notes}, template: {bool(template_path)}")
        
        # Generate outline using LLM with optional speaker notes
        orchestrator = LLMOrchestrator(provider, api_key)
        outline = orchestrator.generate_outline_with_notes(text, guidance, include_notes)
        
        # Create PowerPoint presentation with optional template
        ppt_generator = PPTGenerator(template_path)
        success = ppt_generator.create_presentation(outline)
        
        if not success:
            return jsonify({"error": "Failed to create presentation"}), 500
        
        # Save to temporary file
        temp_dir = tempfile.mkdtemp()
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(temp_dir, filename)
        
        ppt_generator.save_presentation(filepath)
        
        logger.info(f"Presentation generated successfully: {filename}")
        
        # Return file
        return send_file(
            filepath,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        logger.error(f"Error generating presentation: {str(e)}")
        return jsonify({"error": "Internal server error occurred"}), 500

@app.route('/providers', methods=['GET'])
def get_providers():
    """Return list of supported LLM providers"""
    return jsonify({
        "providers": [
            {"id": "openai", "name": "OpenAI (GPT-3.5/GPT-4)", "description": "OpenAI's ChatGPT models"},
            {"id": "anthropic", "name": "Anthropic (Claude)", "description": "Anthropic's Claude models"},
            {"id": "groq", "name": "Groq (Mixtral)", "description": "Groq's fast inference platform"}
        ]
    })

@app.route('/test-template', methods=['POST'])
def test_template():
    """Test endpoint to analyze template layouts"""
    try:
        if 'template' not in request.files:
            return jsonify({"error": "No template file provided"}), 400
        
        template_file = request.files['template']
        if not template_file.filename:
            return jsonify({"error": "No file selected"}), 400
        
        # Save template to temporary location
        temp_dir = tempfile.mkdtemp()
        template_path = os.path.join(temp_dir, "test_template.pptx")
        template_file.save(template_path)
        
        # Analyze template
        presentation = Presentation(template_path)
        
        layouts_info = []
        for i, layout in enumerate(presentation.slide_layouts):
            placeholders_info = []
            
            # Create temporary slide to check placeholders
            temp_slide = presentation.slides.add_slide(layout)
            for placeholder in temp_slide.placeholders:
                placeholders_info.append({
                    "idx": placeholder.placeholder_format.idx,
                    "name": getattr(placeholder, 'name', 'no name'),
                    "type": str(placeholder.placeholder_format.type)
                })
            
            # Remove temporary slide
            rId = presentation.slides._sldIdLst[-1].rId
            presentation.part.drop_rel(rId)
            del presentation.slides._sldIdLst[-1]
            
            layouts_info.append({
                "index": i,
                "name": layout.name,
                "placeholders": placeholders_info
            })
        
        return jsonify({
            "template_name": template_file.filename,
            "total_layouts": len(presentation.slide_layouts),
            "existing_slides": len(presentation.slides),
            "layouts": layouts_info
        })
        
    except Exception as e:
        logger.error(f"Error analyzing template: {str(e)}")
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
